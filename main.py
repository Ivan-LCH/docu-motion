# -----------------------------------------------------------------------------------------------------------------------------#
# 1. Import & Library
# -----------------------------------------------------------------------------------------------------------------------------#
# [OS/시스템] 파일 처리, 비동기, JSON 파싱, 정규식, 로깅
import os, asyncio, json, time, shutil, fitz, re, logging

# [UI] Streamlit 웹 인터페이스
import streamlit as st

# [TTS] Microsoft Edge TTS - 텍스트를 음성으로 변환
import edge_tts

# [유틸] 경로 처리 및 시간
from pathlib import Path
from datetime import datetime

# [이미지] PIL - 이미지 처리 (moviepy 호환성 패치 포함)
from PIL import Image, ImageDraw, ImageFont
if not hasattr(Image, 'ANTIALIAS'):
    Image.ANTIALIAS = Image.LANCZOS  # PIL 10.0+ 버전 호환성

# [PPT] PowerPoint 처리
from pptx import Presentation
import io

# [영상처리] MoviePy - 이미지/오디오를 영상으로 합성
from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips, TextClip, CompositeVideoClip, ColorClip

# [외부API] YouTube 업로드 매니저, Google Gemini AI
import youtube_manager

# [TTS] Qwen3-TTS Manager
try:
    from tts_manager import TTSEngine
except ImportError:
    TTSEngine = None

# -----------------------------------------------------------------------------------------------------------------------------#
# Configuration for Remote TTS Server
# -----------------------------------------------------------------------------------------------------------------------------#
TTS_SERVER_URL = os.getenv("TTS_SERVER_URL", "http://localhost:8002")
TTS_VOICE_NAME = os.getenv("TTS_VOICE_NAME", "myvoice")



# -----------------------------------------------------------------------------------------------------------------------------#
# 2. Logging Setup
# -----------------------------------------------------------------------------------------------------------------------------#
# 모듈별 로거 생성 - INFO 레벨 이상만 기록
logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

# 중복 핸들러 방지: 핸들러가 없을 경우에만 추가
if not logger.handlers:
    stream_handler = logging.StreamHandler()                                         # 콘솔 출력용 핸들러
    file_handler   = logging.FileHandler("app.log", encoding='utf-8')                # 파일 기록용 핸들러 (app.log)
    formatter      = logging.Formatter('%(asctime)s - %(levelname)s - %(message)s')  # 로그 포맷: 시간-레벨-메시지

    stream_handler.setFormatter(formatter)
    file_handler  .setFormatter(formatter)
    logger.addHandler(stream_handler)
    logger.addHandler(file_handler)



# -----------------------------------------------------------------------------------------------------------------------------#
# 3. Configuration & Layout (와이드 레이아웃 유지)
# -----------------------------------------------------------------------------------------------------------------------------#
# [앱 정보] 버전 및 프로젝트명
VERSION        = "3.4.9"
PROJECT_NAME   = "DocuMotion Studio"

st.set_page_config(page_title=PROJECT_NAME, page_icon="🎬", layout="wide")

# [영상 렌더링 설정] 캔버스 크기, 폰트, 자막 스타일
CANVAS_SIZE    = (1280, 720)             # 720p HD 해상도
FONT_PATH      = "font.ttf"              # 자막용 폰트 파일 경로
FONT_SIZE      = 28                      # 자막 폰트 크기 (px)
TEXT_COLOR     = 'white'                 # 자막 텍스트 색상
BG_COLOR       = (0, 0, 0)               # 배경색 (검정)
YT_DESCRIPTION = """AI 기반으로 제작된 자동 생성 영상입니다.

📌 Summary
이 영상은 DocuMotion Studio를 통해 PDF/이미지 자료를 분석하여 제작되었습니다.
핵심 내용 요약과 전문적인 AI 내레이션을 통해 정보를 빠르고 정확하게 습득하세요.

🛠️ Tech Stack
- Guide: Ivan
- Analysis: NotebookLM & Gemini
- TTS: Edge-TTS
- Rendering: MoviePy & DocuMotion Engine

#AI영상 #지식공유 #DocuMotion #자동화 #TechInfo
"""

# [환경 분기] 클라우드(Streamlit Cloud) vs 로컬 환경 구분
IS_CLOUD       = "STREAMLIT_RUNTIME_ENV" in os.environ

# [API 키 로딩] Streamlit Secrets에서 API 키 및 토큰 로드
try:
    GOOGLE_API_KEY = st.secrets["GOOGLE_API_KEY"]
    # 클라우드 환경: secrets에서 YouTube 토큰을 파일로 추출
    if "YOUTUBE_TOKEN_JSON" in st.secrets:
        with open("token.json", "w", encoding="utf-8") as f:
            f.write(st.secrets["YOUTUBE_TOKEN_JSON"])
except Exception as e:
    st.error("🔑 Secrets 설정을 확인하세요."); st.stop()

# [디렉토리 설정] 임시 파일 및 출력 파일 저장 경로
BASE_DIR     = Path(__file__).parent
TEMP_DIR     = BASE_DIR / "temp"          # TTS 오디오, 추출된 이미지 저장
OUTPUT_DIR   = BASE_DIR / "outputs"       # 최종 렌더링된 영상 저장

# 디렉토리 자동 생성
for folder in [TEMP_DIR, OUTPUT_DIR]: folder.mkdir(parents=True, exist_ok=True)



# -----------------------------------------------------------------------------------------------------------------------------#
# 4. Helper Functions
# -----------------------------------------------------------------------------------------------------------------------------#

# ─────────────────────────────────────────────────────────────────────────────
# clear_work_directories: 작업 디렉토리 초기화
# - temp/, outputs/ 폴더 내 모든 파일/폴더 삭제
# - 새 프로젝트 시작 또는 수동 클렌징 시 호출
# ─────────────────────────────────────────────────────────────────────────────
def clear_work_directories():
    for folder in [TEMP_DIR, OUTPUT_DIR]:
        if folder.exists():
            for filename in os.listdir(folder):
                file_path = folder / filename
                try:
                    if os.path.isfile(file_path): os.unlink(file_path)
                    elif os.path.isdir(file_path): shutil.rmtree(file_path)
                except Exception: pass

# ─────────────────────────────────────────────────────────────────────────────
# cleanup_moviepy_temp: MoviePy 임시 파일 정리
# - 루트 디렉토리에 생성되는 TEMP_MPY_wvf_snd.mp3 파일 삭제
# ─────────────────────────────────────────────────────────────────────────────
def cleanup_moviepy_temp():
    import glob
    # TEMP_DIR 내 MoviePy 임시 파일 정리
    for f in glob.glob(str(TEMP_DIR / "*TEMP_MPY*.mp3")):
        try: os.unlink(f)
        except: pass
    # 루트 디렉토리 임시 파일도 정리 (혹시 남아있을 경우)
    for f in glob.glob(str(BASE_DIR / "*TEMP_MPY*.mp3")):
        try: os.unlink(f)
        except: pass

# ─────────────────────────────────────────────────────────────────────────────
# get_video_list: outputs 폴더의 영상 목록 조회
# ─────────────────────────────────────────────────────────────────────────────
def get_video_list():
    videos = []
    if OUTPUT_DIR.exists():
        for f in sorted(OUTPUT_DIR.glob("*.mp4"), key=os.path.getmtime, reverse=True):
            videos.append({
                'path': f,
                'name': f.stem,
                'size': f.stat().st_size / (1024 * 1024),  # MB
                'modified': datetime.fromtimestamp(f.stat().st_mtime).strftime('%Y-%m-%d %H:%M')
            })
    return videos

def delete_video(video_path):
    try:
        os.unlink(video_path)
        # 업로드 상태에서도 삭제
        status = load_upload_status()
        video_name = Path(video_path).stem
        if video_name in status:
            del status[video_name]
            save_upload_status(status)
        st.rerun()
    except Exception as e:
        st.error(f"영상 삭제 실패: {e}")

# ─────────────────────────────────────────────────────────────────────────────
# YouTube 업로드 상태 관리 (JSON 파일 기반)
# ─────────────────────────────────────────────────────────────────────────────
UPLOAD_STATUS_FILE = OUTPUT_DIR / "upload_status.json"

def load_upload_status():
    if UPLOAD_STATUS_FILE.exists():
        try:
            with open(UPLOAD_STATUS_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
        except: pass
    return {}

def save_upload_status(status):
    with open(UPLOAD_STATUS_FILE, "w", encoding="utf-8") as f:
        json.dump(status, f, ensure_ascii=False, indent=2)

def mark_as_uploaded(video_name, youtube_url, video_title=None):
    status = load_upload_status()
    status[video_name] = {
        "uploaded": True,
        "url": youtube_url,
        "title": video_title or video_name,
        "uploaded_at": datetime.now().strftime('%Y-%m-%d %H:%M')
    }
    save_upload_status(status)

def get_upload_status(video_name):
    status = load_upload_status()
    return status.get(video_name, {"uploaded": False})

# ─────────────────────────────────────────────────────────────────────────────
# sanitize_filename: 파일명에 사용 불가한 문자 제거
# ─────────────────────────────────────────────────────────────────────────────
def sanitize_filename(name):
    return re.sub(r'[<>:"/\\|?*]', '_', name)[:50]

# ─────────────────────────────────────────────────────────────────────────────
# split_sentences: 텍스트를 문장 단위로 분리
# - 정규식: 마침표/느낌표/물음표 뒤 공백 기준 분할
# - 자막 타이밍 계산에 사용 (문장별 표시 시간 산출)
# ─────────────────────────────────────────────────────────────────────────────
def split_sentences(text):
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    return [s for s in sentences if s]

# ─────────────────────────────────────────────────────────────────────────────
# create_image_from_text: 텍스트로 단순 슬라이드 이미지 생성 (PPT 대용)
# - LibreOffice 부재로 PPT 렌더링 불가 시 대안
# - 검정 배경에 흰색 텍스트로 내용 표시
# ─────────────────────────────────────────────────────────────────────────────
def create_image_from_text(text, filename, size=CANVAS_SIZE):
    img = Image.new('RGB', size, color=(0, 0, 0))
    d = ImageDraw.Draw(img)
    
    # 폰트 로드 (없으면 기본값)
    try:
        font = ImageFont.truetype(FONT_PATH, 40)
    except:
        font = ImageFont.load_default()
        
    # 텍스트 줄바꿈 처리 (간단한 로직)
    margin = 100
    offset = 100
    for line in text.split('\n'):
        # 너무 긴 줄은 대충 자름 (정교한 wrapping은 생략)
        if len(line) > 50:
             line = line[:50] + "..."
        d.text((margin, offset), line, font=font, fill=(255, 255, 255))
        offset += 60
        
    img.save(filename)

# ─────────────────────────────────────────────────────────────────────────────
# process_pptx: PPT 파일 처리
# - 각 슬라이드의 텍스트 추출 -> 이미지 변환
# ─────────────────────────────────────────────────────────────────────────────
def process_pptx(file_stream, temp_dir):
    prs = Presentation(file_stream)
    assets = []
    
    for i, slide in enumerate(prs.slides):
        # 텍스트 추출
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
        
        slide_text = "\n".join(text_runs).strip()
        if not slide_text: slide_text = f"Slide {i+1} (No Text)"
        
        # 이미지 생성 (TEXT -> PNG)
        target = temp_dir / f"ppt_{i+1:02d}.png"
        create_image_from_text(slide_text, str(target))
        assets.append({'path': target, 'label': f"PPT Slide {i+1}", 'extracted_text': slide_text})
        
    return assets

# ─────────────────────────────────────────────────────────────────────────────
# move_slide: 슬라이드 순서 이동
# ─────────────────────────────────────────────────────────────────────────────
def move_slide(from_idx, to_idx):
    slides = st.session_state.master_slides
    if 0 <= to_idx < len(slides):
        slides[from_idx], slides[to_idx] = slides[to_idx], slides[from_idx]
        st.rerun()

# ─────────────────────────────────────────────────────────────────────────────
# delete_slide: 슬라이드 삭제
# ─────────────────────────────────────────────────────────────────────────────
def delete_slide(idx):
    del st.session_state.master_slides[idx]
    st.rerun()

# ─────────────────────────────────────────────────────────────────────────────
# render_video: 슬라이드 데이터를 영상으로 렌더링
# - 입력: [{"image": Path, "text": str}, ...] 형태의 슬라이드 리스트
# - 처리 흐름:
#   1) 각 슬라이드의 텍스트 → TTS(edge_tts)로 오디오 생성 (ko-KR-SunHiNeural 음성)
#   2) 오디오 길이에 맞춰 이미지 클립 생성
#   3) 문장별 자막 클립 생성 (글자 수 비례로 표시 시간 계산)
#   4) 배경 + 이미지 + 자막 + 오디오 합성
#   5) 모든 슬라이드 연결 후 MP4 파일로 출력
# - 출력: 생성된 영상 파일 경로 (실패 시 None)
# - 진행률: 슬라이드 합성(0-50%) + 영상 인코딩(51-100%)
# ─────────────────────────────────────────────────────────────────────────────
from proglog import ProgressBarLogger

class StreamlitProgressLogger(ProgressBarLogger):
    """MoviePy 인코딩 진행률을 Streamlit 프로그레스바에 연결하는 커스텀 로거"""
    def __init__(self, progress_bar):
        super().__init__()
        self.progress_bar = progress_bar
    
    def bars_callback(self, bar, attr, value, old_value=None):
        # bar='t': 비디오 프레임 처리 진행률
        if bar == 't' and attr == 'index':
            total = self.bars[bar]['total']
            if total > 0:
                # 인코딩 진행률: 50% ~ 100% 구간
                encode_progress = int(value / total * 100)
                overall_progress = 50 + int(encode_progress * 0.5)
                self.progress_bar.progress(overall_progress, f"📼 영상 인코딩 중... ({encode_progress}%)")


# ─────────────────────────────────────────────────────────────────────────────
# load_tts_engine: TTS 엔진 로드 (캐싱 적용)
# ─────────────────────────────────────────────────────────────────────────────
@st.cache_resource
def load_tts_engine():
    if TTSEngine is None:
        return None
    try:
        engine = TTSEngine(server_url=TTS_SERVER_URL, voice_name=TTS_VOICE_NAME)
        return engine
    except Exception as e:
        logger.error(f"TTS Engine Load Error: {e}")
        return None

def render_video(data, video_title="DocuMotion Video"):
    total_slides = len(data)
    final_clips = []
    
    # 통합 프로그레스바 생성
    progress_bar = st.progress(0, text="🚀 렌더링 준비 중...")
    
    # TTS 엔진 로드
    tts_engine = load_tts_engine()
    if not tts_engine:
        st.error("TTS 엔진을 로드할 수 없습니다. requirements.txt를 확인하세요.")
        return None

    try:
        for i, item in enumerate(data):
            if not item['text']: continue
            
            # 슬라이드 진행률: 0% ~ 50% 구간
            slide_progress = int((i / total_slides) * 50)
            progress_bar.progress(slide_progress, f"⏳ 슬라이드 {i+1}/{total_slides} 합성 중... ({slide_progress}%)")
            
            # Step 1: TTS 오디오 생성 (Remote TTS Server)
            a_path = TEMP_DIR / f"v_{i}.wav"
            logger.info(f"Generating TTS via remote server ({TTS_SERVER_URL})")
            
            success = tts_engine.generate(
                text=item['text'],
                output_file=str(a_path)
            )
            
            # TTS 실패 시 Edge-TTS로 폴백
            if not success:
               st.warning(f"슬라이드 {i+1} 오디오 생성 실패. 기본 TTS로 시도합니다.")
               try:
                   asyncio.run(edge_tts.Communicate(item['text'], "ko-KR-SunHiNeural").save(str(a_path)))
               except:
                   pass

            if not os.path.exists(str(a_path)):
                 st.error(f"오디오 파일 생성 실패: {item['text'][:20]}...")
                 continue

            # ─────────────────────────────────────────────────────────────────
            # Step 2: 타이밍 계산 (Pacing & Subtitle Sync)
            # ─────────────────────────────────────────────────────────────────
            a_clip = AudioFileClip(str(a_path))
            
            # 원본 오디오 길이 저장 (set_start 전에 저장해야 정확함)
            original_audio_duration = a_clip.duration
            
            # 슬라이드 앞뒤 여백 (부드러운 전환용)
            PAD_DELAY = 0.5  # 앞뒤 각 0.5초 침묵
            
            # 전체 슬라이드 길이 = 오디오 + 앞뒤 여백
            total_duration = original_audio_duration + (PAD_DELAY * 2)
            
            # 오디오 시작점을 여백 이후로 이동
            a_clip = a_clip.set_start(PAD_DELAY)
            
            # ─────────────────────────────────────────────────────────────────
            # Step 3: 문장 분리 및 자막 타이밍 계산
            # ─────────────────────────────────────────────────────────────────
            sentences = split_sentences(item['text'])
            num_sentences = len(sentences)
            total_chars = sum(len(s) for s in sentences)
            
            # 문장 사이 쉼 (0.15초) - 자연스러운 호흡
            SENTENCE_GAP = 0.15
            total_gap_time = SENTENCE_GAP * (num_sentences - 1) if num_sentences > 1 else 0
            
            # 자막에 사용할 실제 시간 = 원본 오디오 - 문장 사이 쉼
            available_subtitle_time = original_audio_duration - total_gap_time
            
            # ─────────────────────────────────────────────────────────────────
            # Step 4: 배경 및 이미지 클립 생성
            # ─────────────────────────────────────────────────────────────────
            bg_clip = ColorClip(size=CANVAS_SIZE, color=BG_COLOR).set_duration(total_duration)
            img_clip = (
                ImageClip(str(item['image']))
                .resize(height=int(CANVAS_SIZE[1] * 0.85))
                .set_position(('center', 'top'))
                .set_duration(total_duration)
            )
            
            # ─────────────────────────────────────────────────────────────────
            # Step 5: 문장별 자막 클립 생성 (글자 수 비례 + 쉼 타이밍)
            # ─────────────────────────────────────────────────────────────────
            subtitle_clips = []
            # 자막 시작: 오디오 시작 + 약간의 지연 (음성보다 약간 늦게 시작)
            SUBTITLE_DELAY = 0.1  # 자막이 음성보다 0.1초 늦게 시작
            current_start = PAD_DELAY + SUBTITLE_DELAY

            for idx, s in enumerate(sentences):
                # 글자 수 비례로 자막 지속 시간 계산
                char_ratio = len(s) / total_chars if total_chars > 0 else 1.0 / num_sentences
                dur = char_ratio * available_subtitle_time
                
                txt_clip = TextClip(
                    txt       = s, 
                    font      = FONT_PATH, 
                    fontsize  = FONT_SIZE, 
                    color     = TEXT_COLOR, 
                    size      = (CANVAS_SIZE[0] - 100, None), 
                    method    = 'caption', 
                    align     = 'center',
                    interline = 8
                ).set_start(current_start).set_duration(dur).set_position(('center', CANVAS_SIZE[1] - 90))
                
                subtitle_clips.append(txt_clip)
                
                # 다음 자막 시작점: 현재 끝 + 문장 사이 쉼
                current_start += dur
                if idx < num_sentences - 1:
                    current_start += SENTENCE_GAP
                
            # ─────────────────────────────────────────────────────────────────
            # Step 6: 레이어 합성 (배경 → 이미지 → 자막) + 오디오 연결
            # ─────────────────────────────────────────────────────────────────
            final_clips.append(CompositeVideoClip([bg_clip, img_clip] + subtitle_clips).set_audio(a_clip))

        # 슬라이드 합성 완료 (50%)
        progress_bar.progress(50, "📼 영상 인코딩 시작...")
        
        # 모든 슬라이드 연결 및 파일 출력 (커스텀 로거로 진행률 표시)
        safe_title = sanitize_filename(video_title)
        out_path  = OUTPUT_DIR / f"{safe_title}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.mp4"
        st_logger = StreamlitProgressLogger(progress_bar)
        temp_audio_path = str(TEMP_DIR / f"{safe_title}_TEMP_MPY.mp3")
        concatenate_videoclips(final_clips, method="compose").write_videofile(
            str(out_path), fps=24, logger=st_logger, temp_audiofile=temp_audio_path
        )
        
        # MoviePy 임시 파일 정리
        cleanup_moviepy_temp()
        
        # 완료 (100%)
        progress_bar.progress(100, "✅ 렌더링 완료!")
        return out_path

    except Exception as e:
        st.error(f"렌더링 오류: {e}")
        return None


# ─────────────────────────────────────────────────────────────────────────────
# upload_to_youtube: YouTube Shorts 업로드 처리
# - 입력: file_path(영상 경로), title(제목), description(설명)
# - youtube_manager 모듈을 통해 업로드 수행
# - 성공 시 URL 반환, 실패 시 에러 메시지 + 로깅
# - 토큰 만료, 할당량 초과 등의 에러 핸들링 포함
# ─────────────────────────────────────────────────────────────────────────────
def upload_to_youtube(file_path: str, title: str, description: str = "AI Video", tags: str = None):
    """유튜브 업로드 공통 함수"""
    with st.spinner(f"🚀 '{title}' 유튜브 업로드 중..."):
        try:
            url = youtube_manager.upload_short(
                file_path   = file_path,
                title       = title,
                description = description,
                tags        = tags
            )
            if url:
                st.success(f"✅ 업로드 성공: {url}")
                return url
            else:
                st.error("❌ 유튜브 업로드 실패: URL을 반환하지 않았습니다. 인증 또는 할당량을 확인하세요.")
                return None
            
        except Exception as e:
            st.error(f"🚨 유튜브 업로드 중 예외 발생: {str(e)}")
            logger.error(f"YouTube Upload Exception: {e}", exc_info=True)
            return None


# ===================================================================================================================
# Main UI
# ===================================================================================================================
# 앱의 메인 UI 구성
# - 사이드바: 파일 업로드, 유튜브 설정, 클렌징/렌더링 버튼
# - 메인 영역: 타임라인 에디터, JSON 일괄 입력, 영상 미리보기
# - 세션 상태: master_slides(슬라이드 목록), scripts(대사 텍스트), last_v(마지막 렌더링 영상)
# ===================================================================================================================
def main():
    st.title(f"🎬 {PROJECT_NAME} {VERSION}")
    
    # ─────────────────────────────────────────────────────────────
    # 메인 영역 최상단: 생성된 영상 목록
    # ─────────────────────────────────────────────────────────────
    video_list = get_video_list()
    if video_list:
        # 컴팩트 버튼 스타일 CSS (Expander 내부의 버튼만 타겟팅하여 적용)
        st.markdown("""
        <style>
        div[data-testid="stExpanderDetails"] div[data-testid="stButton"] button {
            padding-top: 4px !important;
            padding-bottom: 4px !important;
            padding-left: 10px !important;
            padding-right: 10px !important;
            font-size: 14px !important;
            min-height: 0px !important;
            height: auto !important;
            border: 1px solid #ddd !important;
        }
        div[data-testid="stExpanderDetails"] div[data-testid="stButton"] {
            margin-bottom: 0px !important;
            margin-top: 0px !important;
        }
        </style>
        """, unsafe_allow_html=True)

        with st.expander(f"📁 생성된 영상 목록 ({len(video_list)}개)", expanded=False):
            for v in video_list:
                upload_info = get_upload_status(v['name'])
                is_uploaded = upload_info.get("uploaded", False)
                
                # 업로드된 항목: 제목 표시, 미업로드: 파일명 표시
                display_name = upload_info.get("title", v['name']) if is_uploaded else v['name']
                
                # 컬럼 비율 조정: 제목(6) + 공백(1) + 다운(0.5) + 업로드(0.5) + 삭제(0.5)
                # 오른쪽 끝으로 밀기 위해 앞쪽 컬럼 비중을 높임
                col_title, col_space, col_dl, col_up, col_del = st.columns([5, 1, 0.6, 0.6, 0.6])
                
                with col_title:
                    if is_uploaded:
                        url = upload_info.get("url", "#")
                        # 제목과 링크를 한 줄에 표시
                        st.markdown(f"**{display_name}** &nbsp; [[YouTube]({url})]", unsafe_allow_html=True)
                    else:
                        st.write(f"**{display_name}**")
                    st.caption(f"📅 {v['modified']}")
                
                # 공백 컬럼은 비워둠 (col_title이 넓어서 자연스럽게 밀림)
                
                with col_dl:
                    with open(v['path'], "rb") as f:
                        st.download_button("💾", f, file_name=f"{v['name']}.mp4", key=f"dl_{v['name']}", 
                                          help="다운로드")
                with col_up:
                    if is_uploaded:
                         # 이미 업로드된 경우 버튼 비활성화 대신 체크 표시만 (클릭 불가)
                         st.button("✅", key=f"status_{v['name']}", disabled=True, help="업로드 완료")
                    else:
                        if st.button("📤", key=f"reup_{v['name']}", help="YouTube 업로드"):
                            st.session_state.upload_target_video = str(v['path'])
                            st.session_state.show_upload_dialog = True
                            st.rerun()
                with col_del:
                    if st.button("🗑️", key=f"vdel_{v['name']}", help="삭제"):
                        delete_video(v['path'])
    
    # ─────────────────────────────────────────────────────────────────
    # 사이드바: 파일 업로드 및 설정
    # ─────────────────────────────────────────────────────────────────
    with st.sidebar:
        st.header("📂 교육 자료 업로드")
        
        # 업로드 모드 선택
        upload_mode = st.radio("업로드 모드", ["🔄 기존 교체", "➕ 추가"], horizontal=True, key="upload_mode")
        
        # 추가 모드: 삽입 위치 선택
        insert_position = "맨 뒤"
        position_idx = 0
        if upload_mode == "➕ 추가" and 'master_slides' in st.session_state and len(st.session_state.master_slides) > 0:
            insert_position = st.selectbox("삽입 위치", ["맨 앞", "맨 뒤", "특정 위치"], key="insert_pos")
            if insert_position == "특정 위치":
                position_idx = st.number_input("위치 번호 (1부터 시작)", min_value=1, max_value=len(st.session_state.master_slides)+1, value=1, key="pos_idx") - 1
        
        if IS_CLOUD:
            img_in = st.file_uploader("이미지", type=["png", "jpg"], key="cloud_up")
            uploaded_files = [img_in] if img_in else []
        else:
            uploaded_files = st.file_uploader(
                "자료 업로드 (PDF, 이미지, PPT)", 
                type=["pdf", "png", "jpg", "jpeg", "ppt", "pptx"], 
                accept_multiple_files=True,
                key="unified_up"
            )
        
        
        st.divider()
        # 수동 클렌징: 작업 디렉토리 및 세션 상태 초기화
        if st.button("🧹 수동 클렌징", width='stretch'):
            clear_work_directories()
            if 'master_slides' in st.session_state: del st.session_state.master_slides
            st.rerun()
        render_btn = st.button("🚀 영상 렌더링 시작", type="primary", width='stretch')

    # ─────────────────────────────────────────────────────────────────
    # 메인 영역: 파일 업로드 시 활성화
    # ─────────────────────────────────────────────────────────────────
    if uploaded_files:
        # 파일 처리 함수 (진행률 표시 포함)
        def process_uploaded_files(files, progress_bar):
            new_assets = []
            total_files = len(files)
            
            for idx, up_file in enumerate(files):
                progress_bar.progress((idx / total_files), f"📁 파일 처리 중... ({idx+1}/{total_files})")
                
                # 안전한 파일명 생성 (인덱스 + 타임스탬프)
                safe_base = f"{idx+1:02d}_{datetime.now().strftime('%H%M%S')}"
                
                # 1. PDF 처리
                if up_file.type == "application/pdf":
                    doc = fitz.open(stream=up_file.read(), filetype="pdf")
                    total_pages = len(doc)
                    for i in range(total_pages):
                        progress_bar.progress((idx / total_files) + (i / total_pages / total_files), 
                                              f"📄 PDF 변환 중... {up_file.name} ({i+1}/{total_pages})")
                        target = TEMP_DIR / f"pdf_{safe_base}_p{i+1:02d}.png"
                        doc.load_page(i).get_pixmap(dpi=150).save(str(target))
                        new_assets.append({'path': target, 'label': f"{up_file.name} - P{i+1}", 'script': ""})
                # 2. 이미지 처리
                elif up_file.type.startswith("image/"):
                    ext = Path(up_file.name).suffix or ".png"
                    target = TEMP_DIR / f"img_{safe_base}{ext}"
                    with open(target, "wb") as f: f.write(up_file.getbuffer())
                    new_assets.append({'path': target, 'label': up_file.name, 'script': ""})
                # 3. PPT 처리
                elif up_file.name.endswith(('.ppt', '.pptx')):
                    ppt_assets = process_pptx(up_file, TEMP_DIR)
                    for asset in ppt_assets:
                        asset['script'] = ""  # 통합 구조에 script 추가
                    new_assets.extend(ppt_assets)
            
            progress_bar.progress(1.0, "✅ 파일 처리 완료!")
            return new_assets
        
        # 새 파일 업로드 감지
        if 'current_file_set' not in st.session_state:
            st.session_state.current_file_set = True
            
            # 진행률 표시바 생성
            progress_bar = st.progress(0, text="🚀 파일 처리 준비 중...")
            
            # 업로드 모드에 따른 처리 (파일 처리 전에 디렉토리 정리)
            if upload_mode == "🔄 기존 교체" or 'master_slides' not in st.session_state:
                clear_work_directories()
                new_assets = process_uploaded_files(uploaded_files, progress_bar)
                st.session_state.master_slides = new_assets
            else:  # ➕ 추가 모드
                new_assets = process_uploaded_files(uploaded_files, progress_bar)
                existing = st.session_state.master_slides
                if insert_position == "맨 앞":
                    st.session_state.master_slides = new_assets + existing
                elif insert_position == "맨 뒤":
                    st.session_state.master_slides = existing + new_assets
                else:  # 특정 위치
                    st.session_state.master_slides = existing[:position_idx] + new_assets + existing[position_idx:]
            
            st.rerun()
        
        # 세션 상태 초기화 (최초 1회)
        if 'master_slides' not in st.session_state:
            st.session_state.master_slides = []

        # JSON 일괄 입력: {슬라이드번호: 대사텍스트} 형식의 JSON 파싱
        with st.expander("🛠️ JSON 대사 일괄 입력", expanded=False):
            json_text = st.text_area("JSON 데이터를 붙여넣으세요")
            if st.button("✅ 일괄 적용", width='stretch'):
                try:
                    clean_json = re.sub(r'\[cite.*?\]', '', json_text)  # Gemini 인용 태그 제거
                    data = json.loads(clean_json)
                    for k, v in data.items():
                        idx = int(k)
                        if idx < len(st.session_state.master_slides):
                            st.session_state.master_slides[idx]['script'] = v
                            # 위젯 키도 업데이트 (Streamlit은 key가 있으면 value보다 우선)
                            st.session_state[f"t_{idx}"] = v
                    st.rerun()
                except Exception as e: st.error(f"JSON 오류: {e}")

        # ─────────────────────────────────────────────────────────────
        # 타임라인 에디터: 슬라이드별 이미지 + 대사 입력 + 컨트롤
        # ─────────────────────────────────────────────────────────────
        st.subheader("📑 편집 타임라인")
        for i, slide in enumerate(st.session_state.master_slides):
            with st.container(border=True):
                c1, c2, c3 = st.columns([1, 2, 0.3])
                with c1: 
                    # 파일을 직접 읽어서 표시 (경로 문제 방지)
                    img_path = Path(slide['path'])
                    if img_path.exists():
                        with open(img_path, "rb") as img_file:
                            st.image(img_file.read(), use_container_width=True)
                    else:
                        st.warning(f"이미지 없음: {slide['label']}")
                with c2:
                    # 스크립트 입력 (통합 구조 사용)
                    # 세션 상태에 값이 없으면 기본값 설정
                    if f"t_{i}" not in st.session_state:
                        st.session_state[f"t_{i}"] = slide.get('script', '')
                    new_script = st.text_area(f"Slide {i+1}", key=f"t_{i}", height=120)
                    st.session_state.master_slides[i]['script'] = new_script
                with c3:
                    # 슬라이드 컨트롤 버튼
                    st.write("")  # 간격 맞춤
                    if st.button("⬆️", key=f"up_{i}", disabled=(i == 0)):
                        move_slide(i, i - 1)
                    if st.button("⬇️", key=f"dn_{i}", disabled=(i == len(st.session_state.master_slides) - 1)):
                        move_slide(i, i + 1)
                    if st.button("🗑️", key=f"del_{i}"):
                        delete_slide(i)

        # ─────────────────────────────────────────────────────────────
        # 렌더링 트리거: 영상 생성
        # ─────────────────────────────────────────────────────────────
        if render_btn:
            render_data = [{
                "image" : s['path'], 
                "text"  : s.get('script', '')
            } for s in st.session_state.master_slides]

            video_file = render_video(render_data, "DocuMotion Video")
            if video_file:
                st.session_state.last_v = str(video_file)
                st.video(st.session_state.last_v)

        if 'last_v' in st.session_state:
            st.divider()
            col1, col2 = st.columns(2)
            with col1:
                if st.button("📺 YouTube 업로드", width='stretch'):
                    st.session_state.show_upload_dialog = True
            with col2:
                with open(st.session_state.last_v, "rb") as f:
                    video_name = Path(st.session_state.last_v).stem
                    st.download_button("💾 동영상 다운로드", f, file_name=f"{video_name}.mp4")
        
    else:
        # 파일 미업로드 상태: 세션 정리 및 안내 메시지
        if 'current_file_set' in st.session_state: del st.session_state.current_file_set
        st.info("파일을 업로드하여 시작하세요.")
    
    # ─────────────────────────────────────────────────────────────
    # YouTube 업로드 다이얼로그 (전역)
    # ─────────────────────────────────────────────────────────────
    if st.session_state.get('show_upload_dialog', False):
        # 업로드 대상 영상 결정 (목록에서 선택 또는 마지막 렌더링)
        target_video = st.session_state.get('upload_target_video', st.session_state.get('last_v', ''))
        
        if target_video:
            with st.container(border=True):
                st.subheader("📺 YouTube 업로드 설정")
                default_title = Path(target_video).stem
                video_title = st.text_input("영상 제목", value=default_title, key="yt_title")
                video_desc = st.text_area("영상 설명", value=YT_DESCRIPTION, height=120, key="yt_desc")
                video_tags = st.text_input("태그 (쉼표 구분)", value="DocuMotion, 자동화, AI", key="yt_tags")
                st.caption("📁 카테고리: 과학기술 | 🌐 언어: 한국어")
                
                btn_col1, btn_col2 = st.columns(2)
                with btn_col1:
                    if st.button("✅ 업로드 실행", type="primary", width='stretch'):
                        video_name = Path(target_video).stem
                        url = upload_to_youtube(target_video, video_title, video_desc, video_tags)
                        if url:
                            mark_as_uploaded(video_name, url, video_title)
                            st.success(f"✅ 업로드 완료! [YouTube에서 보기]({url})")
                        st.session_state.show_upload_dialog = False
                        if 'upload_target_video' in st.session_state:
                            del st.session_state.upload_target_video
                        st.rerun()
                with btn_col2:
                    if st.button("❌ 취소", width='stretch'):
                        st.session_state.show_upload_dialog = False
                        if 'upload_target_video' in st.session_state:
                            del st.session_state.upload_target_video
                        st.rerun()


# ===================================================================================================================
# 엔트리 포인트
# ===================================================================================================================
if __name__ == "__main__":
    main()


# ===================================================================================================================
# End of program
# ===================================================================================================================
